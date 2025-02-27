import streamlit as st
import os
from io import BytesIO
from pptx import Presentation
from PIL import Image
import aspose.slides as slides
import aspose.pydrawing as drawing
from aspose.slides import SlideSizeScaleType

# ✅ Function to remove watermark
def remove_watermark(input_pptx: str, output_pptx: str):
    """
    Removes watermark text from the given PowerPoint file.
    :param input_pptx: Path to the input PowerPoint file with watermark.
    :param output_pptx: Path to save the PowerPoint file without watermark.
    """
    prs = Presentation(input_pptx)

    watermark_phrases = [
        "Evaluation only.",
        "Created with Aspose.Slides for Python via .NET 25.1.",
        "Copyright 2004-2025Aspose Pty Ltd.",
        "Copyright 2004-2025 Aspose Pty Ltd."
    ]

    for slide in prs.slides:
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for phrase in watermark_phrases:
                    if phrase in shape.text:
                        shape.text = shape.text.replace(phrase, "").strip()
            if hasattr(shape, "name") and "watermark" in shape.name.lower():
                shapes_to_remove.append(shape)  # Mark for deletion
        
        for shape in shapes_to_remove:
            slide.shapes._spTree.remove(shape._element)

    prs.save(output_pptx)

# ✅ Convert PPT slides to images for display
def convert_ppt_to_images(ppt_file):
    images = []
    with slides.Presentation(ppt_file) as pres:
        for slide in pres.slides:
            image_stream = BytesIO()
            bmp = slide.get_thumbnail(1, 1)
            bmp.save(image_stream, drawing.imaging.ImageFormat.png)
            image_stream.seek(0)
            images.append(image_stream)
    return images

# ✅ Extract slides based on filters
def extract_matching_slides(ppt_file, rationale_filters, keyword):
    prs = Presentation(ppt_file)
    matching_slides = []
    images = convert_ppt_to_images(ppt_file)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]).lower()
        if any(rationale.lower() in slide_text for rationale in rationale_filters) or (keyword.lower() in slide_text if keyword else False):
            matching_slides.append((slide_num, images[slide_num - 1]))
    
    return matching_slides

# ✅ Streamlit UI
st.set_page_config(page_title="Deckorator V2.1", page_icon="📊", layout="wide")

sidebar, main_content = st.columns([8, 2])

with sidebar:
    st.title("Deckorator V2.1")

    with st.container(border=True):
        col1, col2, col3, col4 = st.columns(4)

        with col1:
            sentiment = st.selectbox("Sentiment", ["Positive", "Neutral", "Critic"], index=2)

        with col2:
            rationale_filter = st.multiselect("Rationale", ["Safety/Tolerability", "Efficacy"], default=["Safety/Tolerability"])

        with col3:
            drug_filter = st.multiselect("Drug", ["Mariposa", "Datatumumab", "Teclistamab"], default=[])

        with col4:
            keyword_search = st.text_input("Keyword Search", placeholder="Type to Search")

    backend_ppt_path = "" #ADD PPT PATH HERE
    if not os.path.exists(backend_ppt_path):
        st.error("Add PPT in code / No backend PPT file found!")
    else:
        st.success(f"Loaded PPT : {backend_ppt_path}")

    if "selected_slides" not in st.session_state:
        st.session_state.selected_slides = {}

    if "confirmed" not in st.session_state:
        st.session_state.confirmed = False

    if os.path.exists(backend_ppt_path):
        matching_slides = extract_matching_slides(backend_ppt_path, rationale_filter, keyword_search)
        col10, col11, col12, col13 = st.columns(4)

        with col10:
            st.subheader("Matching Slides")
        with col13:
            select_all = st.checkbox("Select All Slides", key="select_all")

        cols = st.columns(3)
        for idx, (slide_num, slide_image) in enumerate(matching_slides):
            with cols[idx % 3]:
                with st.container():
                    is_selected = st.checkbox("", key=f"Slide_{slide_num}", value=slide_num in st.session_state.selected_slides)
                    st.image(slide_image, caption=f"Slide {slide_num}", use_container_width=True)

                    if is_selected:
                        st.session_state.selected_slides[slide_num] = slide_image
                    elif slide_num in st.session_state.selected_slides:
                        del st.session_state.selected_slides[slide_num]
                    st.write("---")

with main_content:
    st.subheader("Selected Slides")
    with st.container(border=True):
        if st.session_state.selected_slides:
            for slide_num, slide_image in st.session_state.selected_slides.items():
                st.image(slide_image, caption=f"Slide {slide_num}", use_container_width=True)
                st.write("---")

            # Confirm Selection
            if not st.session_state.confirmed:
                if st.button("✔ Confirm Selection"):
                    st.session_state.confirmed = True
                    st.rerun()

            # ✅ Download after Confirmation
            if st.session_state.confirmed:
                with slides.Presentation(backend_ppt_path) as original_ppt:
                    new_ppt = slides.Presentation()

                    # Set slide size to match original
                    new_ppt.slide_size.set_size(
                        original_ppt.slide_size.size.width,
                        original_ppt.slide_size.size.height,
                        SlideSizeScaleType.DO_NOT_SCALE
                    )

                    # Remove default empty slide
                    while len(new_ppt.slides) > 0:
                        new_ppt.slides.remove_at(0)

                    for slide_num in st.session_state.selected_slides:
                        original_slide = original_ppt.slides[slide_num - 1]
                        
                        # ✅ Add a blank slide instead of copying layout
                        new_slide = new_ppt.slides.add_empty_slide(new_ppt.masters[0].layout_slides[0])

                        # ✅ Clone shapes while keeping font color
                        for shape in original_slide.shapes:
                            cloned_shape = new_slide.shapes.add_clone(shape)

                            # 🔹 Preserve font color if shape has text
                            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                                for paragraph in shape.text_frame.paragraphs:
                                    for portion in paragraph.portions:
                                        portion_color = portion.portion_format.fill_format.solid_fill_color.color  # Get font color
                                        for cloned_paragraph in cloned_shape.text_frame.paragraphs:
                                            for cloned_portion in cloned_paragraph.portions:
                                                cloned_portion.portion_format.fill_format.solid_fill_color.color = portion_color  # Set font color

                    # Save the new PPT and remove watermark
                    temp_pptx = "temp_selected_slides.pptx"
                    new_ppt.save(temp_pptx, slides.export.SaveFormat.PPTX)

                    # ✅ Remove watermark before downloading
                    final_pptx = "Selected_Slides.pptx"
                    remove_watermark(temp_pptx, final_pptx)

                    with open(final_pptx, "rb") as f:
                        ppt_bytes = f.read()

                    st.download_button(
                        label="⬇ Download Powerpoint",
                        type="primary",
                        data=ppt_bytes,
                        file_name="Final_Ppt.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
